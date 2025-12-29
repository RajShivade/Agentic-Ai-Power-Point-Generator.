from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
import datetime

# --- Configuration ---
# Define a professional color palette
DARK_BLUE = RGBColor(0, 86, 179)       # Primary heading color, strong accents
ACCENT_BLUE = RGBColor(102, 178, 255)  # Lighter accents, secondary highlights
DARK_GREY = RGBColor(51, 51, 51)       # Main body text color
LIGHT_GREY = RGBColor(150, 150, 150)   # Subtitles, secondary information
WHITE = RGBColor(255, 255, 255)        # Background color

# Define consistent font sizes
TITLE_FONT_SIZE = Pt(48)
SUBTITLE_FONT_SIZE = Pt(30)
HEADING_FONT_SIZE = Pt(36)
BODY_FONT_SIZE = Pt(20)
BULLET_FONT_SIZE = Pt(18)

# Helper function to create standard Title and Content slides
def add_content_slide(prs, title_text, bullet_points):
    """
    Adds a new slide with a title and a list of bullet points.
    """
    slide_layout = prs.slide_layouts[1]  # Layout for Title and Content
    slide = prs.slides.add_slide(slide_layout)

    # Set title properties
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.name = 'Calibri Light'
    title.text_frame.paragraphs[0].font.size = HEADING_FONT_SIZE
    title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE

    # Set content (bullet points) properties
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()  # Clear any default text

    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0  # Main bullet point level
        p.font.name = 'Calibri'
        p.font.size = BODY_FONT_SIZE
        p.font.color.rgb = DARK_GREY
        p.font.bold = False

# Helper function to add sub-bullets to a text frame
def add_sub_bullet_to_text_frame(text_frame, text, level=1):
    """
    Adds a sub-bullet point to a given text frame.
    """
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level  # Indentation level for sub-bullets
    p.font.name = 'Calibri'
    p.font.size = BULLET_FONT_SIZE
    p.font.color.rgb = DARK_GREY
    p.font.bold = False

# --- Presentation Creation ---
prs = Presentation()
prs.slide_width = Inches(10)  # Standard widescreen size
prs.slide_height = Inches(7.5)

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])  # Layout for Title Slide

# Set main title
title = slide.shapes.title
title.text = "Machine Learning: Empowering Intelligence"
title.text_frame.paragraphs[0].font.name = 'Calibri Light'
title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE

# Set subtitle
subtitle = slide.placeholders[1]
subtitle.text = "Unlocking the Power of Data and Algorithms"
subtitle.text_frame.paragraphs[0].font.name = 'Calibri'
subtitle.text_frame.paragraphs[0].font.size = SUBTITLE_FONT_SIZE
subtitle.text_frame.paragraphs[0].font.color.rgb = LIGHT_GREY
subtitle.text_frame.paragraphs[0].alignment = MSO_ANCHOR.BOTTOM

# Add presenter and date using a custom text box for precise positioning
today = datetime.date.today().strftime("%B %d, %Y")
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(8.0), Inches(0.5))
tf = txBox.text_frame

p_presenter = tf.add_paragraph()
p_presenter.text = "Presented by: Senior Python Developer"
p_presenter.font.name = 'Calibri'
p_presenter.font.size = BODY_FONT_SIZE
p_presenter.font.color.rgb = LIGHT_GREY

p_date = tf.add_paragraph()
p_date.text = f"Date: {today}"
p_date.font.name = 'Calibri'
p_date.font.size = BODY_FONT_SIZE
p_date.font.color.rgb = LIGHT_GREY

# Slide 2: Introduction
add_content_slide(prs, "Introduction to Machine Learning", [
    "Machine Learning (ML) is a subset of Artificial Intelligence (AI) that enables systems to learn from data.",
    "Focuses on developing algorithms that can generalize from examples and make predictions or decisions.",
    "Transforms raw data into actionable insights and predictive models with minimal human intervention.",
    "Drives innovation across various industries by automating tasks and revealing hidden patterns."
])

# Slide 3: Problem Statement / Overview
add_content_slide(prs, "Why Machine Learning Matters", [
    "Solving Complex Problems: Addresses challenges too intricate for traditional programming rules.",
    "Handling Big Data: Efficiently processes and extracts valuable information from massive datasets.",
    "Automation & Efficiency: Automates repetitive tasks, optimizes workflows, and reduces manual effort.",
    "Predictive Capabilities: Forecasts future trends, behaviors, and outcomes with high accuracy.",
    "Continuous Improvement: Models can adapt and enhance performance over time as more data becomes available."
])

# Slide 4: Key Concepts (Types of ML)
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# Set title
title = slide.shapes.title
title.text = "Core Paradigms of Machine Learning"
title.text_frame.paragraphs[0].font.name = 'Calibri Light'
title.text_frame.paragraphs[0].font.size = HEADING_FONT_SIZE
title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE

# Set content with main bullets and sub-bullets
body = slide.shapes.placeholders[1]
tf = body.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Supervised Learning: Learns from labeled data to make predictions."
p.level = 0
p.font.name = 'Calibri'
p.font.size = BODY_FONT_SIZE
p.font.color.rgb = DARK_GREY
add_sub_bullet_to_text_frame(tf, "Classification (e.g., spam detection, image recognition)")
add_sub_bullet_to_text_frame(tf, "Regression (e.g., house price prediction, stock market forecasting)")

p = tf.add_paragraph()
p.text = "Unsupervised Learning: Discovers hidden patterns and structures in unlabeled data."
p.level = 0
p.font.name = 'Calibri'
p.font.size = BODY_FONT_SIZE
p.font.color.rgb = DARK_GREY
add_sub_bullet_to_text_frame(tf, "Clustering (e.g., customer segmentation, anomaly detection)")
add_sub_bullet_to_text_frame(tf, "Dimensionality Reduction (e.g., data visualization, noise reduction)")

p = tf.add_paragraph()
p.text = "Reinforcement Learning: Learns through trial and error by interacting with an environment to maximize rewards."
p.level = 0
p.font.name = 'Calibri'
p.font.size = BODY_FONT_SIZE
p.font.color.rgb = DARK_GREY
add_sub_bullet_to_text_frame(tf, "Applications in robotics, game AI, autonomous systems.")

# Slide 5: Machine Learning Workflow
add_content_slide(prs, "The Machine Learning Lifecycle", [
    "1. Data Collection: Gathering relevant and high-quality raw data from various sources.",
    "2. Data Preprocessing: Cleaning, transforming, and preparing data for model training (e.g., handling missing values, feature scaling).",
    "3. Model Selection: Choosing the appropriate algorithm or model architecture based on the problem type and data characteristics.",
    "4. Model Training: Feeding the preprocessed data to the chosen algorithm to learn patterns and relationships.",
    "5. Model Evaluation: Assessing the performance and accuracy of the trained model using unseen test data.",
    "6. Model Deployment: Integrating the validated model into production systems for real-world use.",
    "7. Monitoring & Retraining: Continuously observing model performance and updating/retraining it as data distributions shift."
])

# Slide 6: Tools & Technologies
add_content_slide(prs, "Essential ML Tools & Technologies", [
    "Programming Languages: Python (dominant), R, Java, Scala.",
    "Core Libraries (Python): NumPy, Pandas (data manipulation), Matplotlib, Seaborn (data visualization).",
    "ML Frameworks: Scikit-learn (traditional ML), XGBoost, LightGBM (gradient boosting).",
    "Deep Learning Frameworks: TensorFlow, PyTorch, Keras (for neural networks).",
    "Cloud Platforms: AWS (SageMaker), Google Cloud (AI Platform), Azure ML (managed ML services).",
    "Big Data Tools: Apache Spark, Hadoop (for large-scale data processing)."
])

# Slide 7: Applications / Use Cases
add_content_slide(prs, "Real-World Applications of Machine Learning", [
    "Healthcare: Disease diagnosis, drug discovery, personalized treatment plans.",
    "Finance: Fraud detection, algorithmic trading, credit risk assessment.",
    "E-commerce: Product recommendation systems, personalized shopping experiences, demand forecasting.",
    "Autonomous Vehicles: Object detection, path planning, predictive maintenance.",
    "Natural Language Processing (NLP): Sentiment analysis, chatbots, machine translation.",
    "Computer Vision: Facial recognition, image classification, quality control in manufacturing."
])

# Slide 8: Benefits / Impact
add_content_slide(prs, "The Transformative Impact of ML", [
    "Enhanced Decision Making: Provides data-driven insights for better strategic and operational choices.",
    "Increased Efficiency: Automates complex and repetitive tasks, freeing up human resources.",
    "Personalization: Delivers highly tailored experiences for customers and users, improving satisfaction.",
    "Innovation & Discovery: Accelerates research, uncovers new patterns, and fosters groundbreaking advancements.",
    "Cost Reduction: Optimizes resource allocation, supply chains, and operational processes.",
    "Competitive Advantage: Drives growth, market leadership, and the ability to adapt quickly to changes."
])

# Slide 9: Challenges / Limitations
add_content_slide(prs, "Challenges and Ethical Considerations in ML", [
    "Data Quality & Bias: Performance is heavily dependent on clean, representative data; biased data leads to biased models.",
    "Model Interpretability (Explainability): Understanding why a 'black box' model makes certain predictions can be difficult, especially in critical applications.",
    "Overfitting & Underfitting: Models performing poorly on unseen data due to learning noise or being too simplistic.",
    "Scalability: Managing and processing ever-growing datasets and deploying models at scale requires significant infrastructure.",
    "Ethical & Privacy Concerns: Misuse of data, algorithmic discrimination, data security, and compliance with regulations.",
    "Resource Intensive: Requires significant computational power, specialized expertise, and ongoing maintenance."
])

# Slide 10: Conclusion
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# Set title
title = slide.shapes.title
title.text = "Conclusion & Future Outlook"
title.text_frame.paragraphs[0].font.name = 'Calibri Light'
title.text_frame.paragraphs[0].font.size = HEADING_FONT_SIZE
title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE

# Set content
body = slide.shapes.placeholders[1]
tf = body.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "ML is a Game-Changer: Revolutionizing industries and daily life, becoming an indispensable tool."
p.level = 0
p.font.name = 'Calibri'
p.font.size = BODY_FONT_SIZE
p.font.color.rgb = DARK_GREY

p = tf.add_paragraph()
p.text = "Continuous Evolution: Expect rapid advancements in algorithms, hardware, and data availability."
p.level = 0
p.font.name = 'Calibri'
p.font.size = BODY_FONT_SIZE
p.font.color.rgb = DARK_GREY

p = tf.add_paragraph()
p.text = "Ethical AI is Paramount: Responsible development, transparent models, and fair deployment are crucial for societal trust."
p.level = 0
p.font.name = 'Calibri'
p.font.size = BODY_FONT_SIZE
p.font.color.rgb = DARK_GREY

p = tf.add_paragraph()
p.text = "Future Trends: Emergence of AutoML, Explainable AI (XAI), Edge AI, and Quantum Machine Learning."
p.level = 0
p.font.name = 'Calibri'
p.font.size = BODY_FONT_SIZE
p.font.color.rgb = DARK_GREY

p = tf.add_paragraph()
p.text = "Embrace the Future: Harness ML for innovation, efficiency, and solving humanity's grand challenges."
p.level = 0
p.font.name = 'Calibri'
p.font.size = BODY_FONT_SIZE
p.font.color.rgb = DARK_BLUE # Emphasize this point

p = tf.add_paragraph()
p.text = "Thank you for your attention!"
p.level = 0
p.font.name = 'Calibri Light'
p.font.size = BODY_FONT_SIZE
p.font.color.rgb = DARK_GREY

p = tf.add_paragraph()
p.text = "Questions?"
p.level = 0
p.font.name = 'Calibri Light'
p.font.size = HEADING_FONT_SIZE # Prominent for Q&A
p.font.color.rgb = DARK_BLUE
p.alignment = MSO_ANCHOR.BOTTOM

# Save the presentation
filename = "Machine_Learning_Presentation.pptx"
prs.save(filename)